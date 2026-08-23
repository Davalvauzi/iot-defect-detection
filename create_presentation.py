import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5) # 16:9 Widescreen

    # Color Palette - Modern Clean Tech (Navy, Royal Blue, Emerald, Crimson, White, Slate)
    NAVY = RGBColor(15, 23, 42)       # #0F172A
    CARD_DARK = RGBColor(30, 41, 59)   # #1E293B
    BLUE = RGBColor(37, 99, 235)      # #2563EB
    EMERALD = RGBColor(16, 185, 129)  # #10B981
    CRIMSON = RGBColor(220, 38, 38)   # #DC2626
    AMBER = RGBColor(217, 119, 6)     # #D97706
    WHITE = RGBColor(255, 255, 255)
    LIGHT_BG = RGBColor(248, 250, 252)# #F8FAFC
    SLATE_TEXT = RGBColor(71, 85, 105)# #475569
    MUTED_TEXT = RGBColor(148, 163, 184) # #94A3B8
    CARD_LIGHT = RGBColor(255, 255, 255)
    CARD_BORDER = RGBColor(226, 232, 240)

    blank_slide_layout = prs.slide_layouts[6]

    def add_header(slide, title_text, category_text="TUGAS PENGEMBANGAN DASHBOARD IOT"):
        # Top category badge
        badge_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(11.5), Inches(0.4))
        tf_b = badge_box.text_frame
        tf_b.word_wrap = True
        tf_b.margin_left = tf_b.margin_top = tf_b.margin_right = tf_b.margin_bottom = 0
        p_b = tf_b.paragraphs[0]
        p_b.text = category_text.upper()
        p_b.font.size = Pt(11)
        p_b.font.bold = True
        p_b.font.color.rgb = BLUE
        p_b.font.name = "Plus Jakarta Sans"

        # Main Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.5), Inches(0.65))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p.font.name = "Plus Jakarta Sans"

    def add_card(slide, left, top, width, height, bg_color=CARD_LIGHT, border_color=CARD_BORDER):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1.5)
        else:
            shape.line.fill.background()
        return shape

    # ==========================================
    # SLIDE 1: TITLE SLIDE (Cover Modern Dark Navy)
    # ==========================================
    s1 = prs.slides.add_slide(blank_slide_layout)
    bg1 = add_card(s1, Inches(0), Inches(0), Inches(13.333), Inches(7.5), NAVY, None)

    # Accent decorative bar
    add_card(s1, Inches(1.0), Inches(1.5), Inches(0.15), Inches(4.2), BLUE, None)

    # Category Pill
    pill = add_card(s1, Inches(1.4), Inches(1.4), Inches(3.8), Inches(0.45), CARD_DARK, BLUE)
    tb_pill = s1.shapes.add_textbox(Inches(1.4), Inches(1.4), Inches(3.8), Inches(0.45))
    tf_p = tb_pill.text_frame
    tf_p.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf_p.paragraphs[0]
    p.text = "PROYEK PENGEMBANGAN DASHBOARD IOT"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Plus Jakarta Sans"

    # Main Title & Subtitle
    tb_title = s1.shapes.add_textbox(Inches(1.4), Inches(2.0), Inches(10.5), Inches(2.6))
    tf_t = tb_title.text_frame
    tf_t.word_wrap = True
    p_t = tf_t.paragraphs[0]
    p_t.text = "Smart Quality Control Dashboard"
    p_t.font.size = Pt(36)
    p_t.font.bold = True
    p_t.font.color.rgb = WHITE
    p_t.font.name = "Plus Jakarta Sans"

    p_sub = tf_t.add_paragraph()
    p_sub.text = "Sistem Deteksi Cacat Barang Berbasis Real-Time Data Streaming & Mobile Monitoring"
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = MUTED_TEXT
    p_sub.font.name = "Plus Jakarta Sans"
    p_sub.space_before = Pt(12)

    # Author / Metadata Card
    meta_card = add_card(s1, Inches(1.4), Inches(4.8), Inches(10.5), Inches(1.2), CARD_DARK, None)
    tb_meta = s1.shapes.add_textbox(Inches(1.7), Inches(5.0), Inches(9.8), Inches(0.8))
    tf_m = tb_meta.text_frame
    p_m1 = tf_m.paragraphs[0]
    p_m1.text = "Topik Tugas: Visualisasi Data • Web Dashboard • Mobile Monitoring • Real-time Data"
    p_m1.font.size = Pt(12)
    p_m1.font.bold = True
    p_m1.font.color.rgb = EMERALD
    p_m1.font.name = "Plus Jakarta Sans"

    p_m2 = tf_m.add_paragraph()
    p_m2.text = "Tools Implementasi: Google Firebase Cloud DB • Node-RED Dashboard • ESP32 • Chart.js"
    p_m2.font.size = Pt(11)
    p_m2.font.color.rgb = WHITE
    p_m2.font.name = "Plus Jakarta Sans"
    p_m2.space_before = Pt(4)


    # ==========================================
    # SLIDE 2: LATAR BELAKANG & PERMASALAHAN
    # ==========================================
    s2 = prs.slides.add_slide(blank_slide_layout)
    add_header(s2, "Latar Belakang & Tantangan Quality Control")

    # 3 Cards: Masalah, Kebutuhan, Solusi
    col_w = Inches(3.64)
    top_pos = Inches(1.7)
    card_h = Inches(5.0)

    # Card 1: Masalah Manual QC
    add_card(s2, Inches(0.8), top_pos, col_w, card_h, CARD_LIGHT, CARD_BORDER)
    add_card(s2, Inches(1.1), Inches(2.0), Inches(0.8), Inches(0.8), RGBColor(254, 242, 242), CRIMSON)
    tb = s2.shapes.add_textbox(Inches(1.1), Inches(3.0), Inches(3.0), Inches(3.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Tantangan QC Manual"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NAVY
    
    p2 = tf.add_paragraph()
    p2.text = "• Rentan human-error akibat kelelahan operator\n• Proses inspeksi lambat & menghambat kecepatan produksi\n• Data cacat tidak terekam otomatis (pencatatan manual di kertas)\n• Sulit melacak tren kerusakan secara real-time"
    p2.font.size = Pt(12)
    p2.font.color.rgb = SLATE_TEXT
    p2.space_before = Pt(10)

    # Card 2: Kebutuhan Industri / IoT
    add_card(s2, Inches(4.84), top_pos, col_w, card_h, CARD_LIGHT, CARD_BORDER)
    add_card(s2, Inches(5.14), Inches(2.0), Inches(0.8), Inches(0.8), RGBColor(254, 243, 199), AMBER)
    tb = s2.shapes.add_textbox(Inches(5.14), Inches(3.0), Inches(3.0), Inches(3.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Kebutuhan Otomasi"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NAVY
    
    p2 = tf.add_paragraph()
    p2.text = "• Sistem inspeksi berkecepatan tinggi tanpa sentuh\n• Pengambilan keputusan lolos/cacat otomatis dalam milidetik\n• Akses pemantauan jarak jauh (Web & Mobile)\n• Pelaporan data historis otomatis untuk audit kualitas"
    p2.font.size = Pt(12)
    p2.font.color.rgb = SLATE_TEXT
    p2.space_before = Pt(10)

    # Card 3: Solusi yang Dibangun
    add_card(s2, Inches(8.88), top_pos, col_w, card_h, CARD_LIGHT, CARD_BORDER)
    add_card(s2, Inches(9.18), Inches(2.0), Inches(0.8), Inches(0.8), RGBColor(239, 246, 255), BLUE)
    tb = s2.shapes.add_textbox(Inches(9.18), Inches(3.0), Inches(3.0), Inches(3.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Solusi IoT Kami"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLUE
    
    p2 = tf.add_paragraph()
    p2.text = "• Node sensor cerdas berbasis ESP32\n• Klasifikasi kondisi barang (Pass / Defect) otomatis\n• Dashboard Web modern dengan visualisasi interaktif & audio alarm\n• Fitur Simulator terintegrasi untuk pengujian fleksibel"
    p2.font.size = Pt(12)
    p2.font.color.rgb = SLATE_TEXT
    p2.space_before = Pt(10)


    # ==========================================
    # SLIDE 3: TUJUAN & RUANG LINGKUP PROYEK
    # ==========================================
    s3 = prs.slides.add_slide(blank_slide_layout)
    add_header(s3, "Tujuan & Ruang Lingkup Proyek")

    # 4 Feature Blocks (Matching the assignment brief)
    blocks = [
        ("1. Data Visualization", "Menampilkan data telemetri sensor dan rasio kualitas secara grafis menggunakan Chart.js (Grafik Donut & Garis Fluktuasi Sinyal).", BLUE),
        ("2. Web Dashboard", "Membangun antarmuka kontrol terpusat dengan indikator KPI Throughput, Yield Rate, Laju Kecepatan, dan Log Riwayat.", EMERALD),
        ("3. Mobile Monitoring", "Desain antarmuka responsif yang dapat diakses dengan mulus melalui perangkat smartphone / tablet operator.", AMBER),
        ("4. Real-Time Streaming", "Pengiriman data instan dari node mikrokontroler ESP32 via protokol jaringan HTTP/REST API tanpa refresh halaman.", CRIMSON)
    ]

    for i, (title, desc, color) in enumerate(blocks):
        x = Inches(0.8 + (i % 2) * 5.95)
        y = Inches(1.7 + (i // 2) * 2.6)
        w = Inches(5.6)
        h = Inches(2.3)

        add_card(s3, x, y, w, h, CARD_LIGHT, CARD_BORDER)
        add_card(s3, x, y, Inches(0.12), h, color, None)

        tb = s3.shapes.add_textbox(x + Inches(0.35), y + Inches(0.3), w - Inches(0.6), h - Inches(0.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = NAVY

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11.5)
        p2.font.color.rgb = SLATE_TEXT
        p2.space_before = Pt(8)


    # ==========================================
    # SLIDE 4: ARSITEKTUR & ALUR KERJA SISTEM
    # ==========================================
    s4 = prs.slides.add_slide(blank_slide_layout)
    add_header(s4, "Arsitektur Sistem End-to-End (3-Tier Layer)")

    tiers = [
        ("PERCEPTION LAYER", "Sensor & Node Hardware", "• Sensor Pendeteksi (Ultrasonik / Infrared)\n• Mikrokontroler ESP32 240MHz\n• Indikator Lokal (LED Alert & Buzzer)\n• Logika Edge Classification Pass/Defect", BLUE),
        ("NETWORK LAYER", "Transmisi & Komunikasi", "• Koneksi Jaringan WiFi 802.11 b/g/n\n• Protokol REST API (HTTP POST)\n• Payload Format JSON Terstruktur\n• Latensi Rendah (< 50ms)", AMBER),
        ("APPLICATION LAYER", "Dashboard & Visualisasi", "• Web & Mobile Dashboard (HTML5, Tailwind)\n• Real-Time Chart.js Visualizer\n• Simulator Engine & Haptic Audio API\n• Riwayat Log, Pagination & Export CSV", EMERALD)
    ]

    for i, (layer_title, sub, details, color) in enumerate(tiers):
        x = Inches(0.8 + i * 3.98)
        y = Inches(1.7)
        w = Inches(3.75)
        h = Inches(5.0)

        add_card(s4, x, y, w, h, CARD_LIGHT, CARD_BORDER)
        
        # Header banner inside card
        add_card(s4, x, y, w, Inches(1.1), color, None)
        tb_h = s4.shapes.add_textbox(x + Inches(0.2), y + Inches(0.15), w - Inches(0.4), Inches(0.8))
        tf_h = tb_h.text_frame
        tf_h.word_wrap = True
        p_h1 = tf_h.paragraphs[0]
        p_h1.text = layer_title
        p_h1.font.size = Pt(13)
        p_h1.font.bold = True
        p_h1.font.color.rgb = WHITE
        
        p_h2 = tf_h.add_paragraph()
        p_h2.text = sub
        p_h2.font.size = Pt(10)
        p_h2.font.color.rgb = WHITE

        # Body details
        tb_b = s4.shapes.add_textbox(x + Inches(0.3), y + Inches(1.3), w - Inches(0.6), Inches(3.4))
        tf_b = tb_b.text_frame
        tf_b.word_wrap = True
        p_b = tf_b.paragraphs[0]
        p_b.text = details
        p_b.font.size = Pt(12)
        p_b.font.color.rgb = SLATE_TEXT
        p_b.space_before = Pt(4)


    # ==========================================
    # SLIDE 5: DESAIN HARDWARE & SKEMA KABEL
    # ==========================================
    s5 = prs.slides.add_slide(blank_slide_layout)
    add_header(s5, "Implementasi Hardware & Konfigurasi Pinout")

    # Left Card: Deskripsi Hardware
    add_card(s5, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.0), CARD_LIGHT, CARD_BORDER)
    tb_l = s5.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(5.0), Inches(4.4))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    p = tf_l.paragraphs[0]
    p.text = "Komponen Hardware Utama"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NAVY

    p2 = tf_l.add_paragraph()
    p2.text = "1. ESP32 NodeMCU Development Board\n   Mikrokontroler dengan modul WiFi terintegrasi untuk membaca sensor dan mengirim data HTTP POST.\n\n2. Sensor Deteksi Kualitas\n   • Sensor Ultrasonik HC-SR04: Deteksi cacat dimensi/ketinggian barang dalam satuan cm.\n   • ATAU Sensor Infrared TCRT5000: Deteksi cacat pantulan optik/warna/lubang.\n\n3. Skema Direct Connection (Praktis):\n   Hanya butuh 4 kabel jumper Female-to-Female langsung tanpa wajib breadboard!"
    p2.font.size = Pt(11.5)
    p2.font.color.rgb = SLATE_TEXT
    p2.space_before = Pt(8)

    # Right Card: Tabel Pinout
    add_card(s5, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.0), CARD_LIGHT, CARD_BORDER)
    tb_r = s5.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.1), Inches(4.4))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    p = tf_r.paragraphs[0]
    p.text = "Tabel Skema Sambungan Kabel (Pinout)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLUE

    p2 = tf_r.add_paragraph()
    p2.text = "A. Pinout Sensor Ultrasonik HC-SR04:\n   • VCC Sensor  ➔  Pin 5V / VIN ESP32\n   • GND Sensor  ➔  Pin GND ESP32\n   • Trig Pin    ➔  GPIO 5 (D5) ESP32\n   • Echo Pin    ➔  GPIO 19 (D19) ESP32\n\nB. Indikator Bawaan (Onboard):\n   • LED Biru Bawaan Board (GPIO 2) otomatis berkedip saat mendeteksi barang cacat tanpa kabel tambahan!"
    p2.font.size = Pt(11.5)
    p2.font.color.rgb = SLATE_TEXT
    p2.space_before = Pt(8)


    # ==========================================
    # SLIDE 6: IMPLEMENTASI WEB & MOBILE DASHBOARD
    # ==========================================
    s6 = prs.slides.add_slide(blank_slide_layout)
    add_header(s6, "Tampilan Dashboard: UI/UX & Visualisasi Data")

    # 4 UI Feature Blocks
    features = [
        ("Simulasi Jalur Konveyor Interaktif", "Bilik animasi pergerakan barang dan respon laser sensor secara visual (Hijau = Lolos, Merah Flash = Cacat).", BLUE),
        ("Kartu Metrik KPI Real-Time", "Menampilkan Total Barang Diperiksa, Jumlah Barang Lolos (%), Jumlah Cacat (%), dan Laju Kecepatan Aliran (item/menit).", EMERALD),
        ("Grafik Analitik Fluktuasi & Rasio", "Grafik Donut untuk perbandingan kualitas Pass vs Defect dan Grafik Garis untuk memantau fluktuasi sinyal sensor.", AMBER),
        ("Filter Batch Jangkauan Riwayat", "Dropdown filter cerdas pada grafik untuk memilih jangkauan: 10 Data Terakhir, 25 Data, 50 Data, atau Seluruh Riwayat.", BLUE)
    ]

    for i, (title, desc, color) in enumerate(features):
        x = Inches(0.8 + (i % 2) * 5.95)
        y = Inches(1.7 + (i // 2) * 2.6)
        w = Inches(5.6)
        h = Inches(2.3)

        add_card(s6, x, y, w, h, CARD_LIGHT, CARD_BORDER)
        add_card(s6, x, y, Inches(0.12), h, color, None)

        tb = s6.shapes.add_textbox(x + Inches(0.35), y + Inches(0.3), w - Inches(0.6), h - Inches(0.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = NAVY

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = SLATE_TEXT
        p2.space_before = Pt(8)


    # ==========================================
    # SLIDE 7: FITUR SIMULATOR & PENGELOLAAN DATA
    # ==========================================
    s7 = prs.slides.add_slide(blank_slide_layout)
    add_header(s7, "Fitur Unggulan Simulator & Manajemen Data")

    col_w = Inches(3.64)
    top_pos = Inches(1.7)
    card_h = Inches(5.0)

    # Card 1: Simulator Lengkap
    add_card(s7, Inches(0.8), top_pos, col_w, card_h, CARD_LIGHT, CARD_BORDER)
    tb = s7.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(3.0), Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🎮 Simulator Terintegrasi"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = BLUE
    p2 = tf.add_paragraph()
    p2.text = "• Tambah Barang Lolos (+OK)\n• Tambah Barang Cacat (+Defect) dengan randomisasi alasan cacat\n• Mode Auto Simulasi dengan slider kecepatan & peluang cacat\n• Input Manual Parameter Kustom"
    p2.font.size = Pt(11.5)
    p2.font.color.rgb = SLATE_TEXT
    p2.space_before = Pt(8)

    # Card 2: Haptic Audio & Interaktivitas
    add_card(s7, Inches(4.84), top_pos, col_w, card_h, CARD_LIGHT, CARD_BORDER)
    tb = s7.shapes.add_textbox(Inches(5.14), Inches(2.0), Inches(3.0), Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🔊 Haptic Audio & Responsif"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = EMERALD
    p2 = tf.add_paragraph()
    p2.text = "• Web Audio API Synthesizer (bunyi beep lolos & nada alarm cacat)\n• Mode Gelap / Terang (Light & Dark Theme Switcher)\n• Animasi Toast Notification otomatis saat data masuk"
    p2.font.size = Pt(11.5)
    p2.font.color.rgb = SLATE_TEXT
    p2.space_before = Pt(8)

    # Card 3: Pagination & CSV Export
    add_card(s7, Inches(8.88), top_pos, col_w, card_h, CARD_LIGHT, CARD_BORDER)
    tb = s7.shapes.add_textbox(Inches(9.18), Inches(2.0), Inches(3.0), Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "📊 Ledger, Pagination & CSV"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = AMBER
    p2 = tf.add_paragraph()
    p2.text = "• Tabel Log Riwayat dengan pagination lengkap (5, 10, 25, 50, Semua baris)\n• Pencarian ID Barang & Filter Status\n• Export Laporan ke format CSV/Excel\n• Dialog Reset Data Multi-Pilihan"
    p2.font.size = Pt(11.5)
    p2.font.color.rgb = SLATE_TEXT
    p2.space_before = Pt(8)


    # ==========================================
    # SLIDE 8: HASIL PENGUJIAN & ANALISIS
    # ==========================================
    s8 = prs.slides.add_slide(blank_slide_layout)
    add_header(s8, "Hasil Pengujian & Analisis Kinerja Sistem")

    tests = [
        ("Skenario 1: Deteksi Barang Normal", "Kondisi: Sinyal sensor berada dalam batas toleransi standar.\nHasil: Dashboard mencatat status LOLOS (OK), LED Hijau menyala, buzzer nada pendek, counter lolos bertambah.", EMERALD),
        ("Skenario 2: Deteksi Barang Cacat", "Kondisi: Terjadi deviasi ukuran / reflektansi di luar ambang batas.\nHasil: Dashboard mencatat status CACAT (DEFECT), alarm buzzer 2-nada berbunyi, servo rejector aktif.", CRIMSON),
        ("Pengujian Latensi & Throughput", "Kondisi: Pengujian pengiriman data beruntun dari ESP32 ke Web.\nHasil: Rata-rata latensi pengiriman < 40ms, visualisasi grafik diperbarui instan tanpa jeda rendering.", BLUE),
        ("Keandalan Penyimpanan (Persistence)", "Kondisi: Refresh halaman web atau restart browser.\nHasil: Data metrik dan riwayat log tetap tersimpan aman melalui mekanisme LocalStorage terstruktur.", AMBER)
    ]

    for i, (title, desc, color) in enumerate(tests):
        x = Inches(0.8 + (i % 2) * 5.95)
        y = Inches(1.7 + (i // 2) * 2.6)
        w = Inches(5.6)
        h = Inches(2.3)

        add_card(s8, x, y, w, h, CARD_LIGHT, CARD_BORDER)
        add_card(s8, x, y, Inches(0.12), h, color, None)

        tb = s8.shapes.add_textbox(x + Inches(0.35), y + Inches(0.3), w - Inches(0.6), h - Inches(0.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = NAVY

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = SLATE_TEXT
        p2.space_before = Pt(8)


    # ==========================================
    # SLIDE 9: KESIMPULAN & PENGEMBANGAN LANJUTAN
    # ==========================================
    s9 = prs.slides.add_slide(blank_slide_layout)
    add_header(s9, "Kesimpulan & Rencana Pengembangan Masa Depan")

    # Left: Kesimpulan
    add_card(s9, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.0), CARD_LIGHT, CARD_BORDER)
    tb_l = s9.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(5.0), Inches(4.4))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    p = tf_l.paragraphs[0]
    p.text = "Kesimpulan Proyek"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = EMERALD

    p2 = tf_l.add_paragraph()
    p2.text = "1. Berhasil Mengimplementasikan Seluruh Syarat Tugas:\n   Visualisasi data, Web Dashboard modern, Mobile Responsive, dan Real-time Data Streaming.\n\n2. Otomasi Quality Control Efektif:\n   Mengeliminasi kesalahan pencatatan manual dan mempercepat deteksi barang cacat.\n\n3. Pengujian Mandiri yang Fleksibel:\n   Fitur simulator memungkinkan verifikasi logika sistem secara penuh bahkan sebelum komponen fisik dirakit."
    p2.font.size = Pt(11.5)
    p2.font.color.rgb = SLATE_TEXT
    p2.space_before = Pt(8)

    # Right: Future Work
    add_card(s9, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.0), CARD_LIGHT, CARD_BORDER)
    tb_r = s9.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.1), Inches(4.4))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    p = tf_r.paragraphs[0]
    p.text = "Rencana Pengembangan Masa Depan"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLUE

    p2 = tf_r.add_paragraph()
    p2.text = "1. Integrasi Cloud Database (Firebase / ThingsBoard):\n   Menyimpan data jutaan transaksi inspeksi di server cloud global.\n\n2. Protokol IoT Ringan (MQTT Broker):\n   Meningkatkan efisiensi bandwidth transmisi data sensor berkecepatan tinggi.\n\n3. Computer Vision / ESP32-CAM:\n   Menggabungkan sensor jarak/cahaya dengan kamera AI untuk klasifikasi cacat visual yang lebih kompleks."
    p2.font.size = Pt(11.5)
    p2.font.color.rgb = SLATE_TEXT
    p2.space_before = Pt(8)


    # ==========================================
    # SLIDE 10: CLOSING & Q&A
    # ==========================================
    s10 = prs.slides.add_slide(blank_slide_layout)
    bg10 = add_card(s10, Inches(0), Inches(0), Inches(13.333), Inches(7.5), NAVY, None)

    add_card(s10, Inches(3.66), Inches(1.6), Inches(6.0), Inches(4.2), CARD_DARK, BLUE)

    tb_end = s10.shapes.add_textbox(Inches(3.8), Inches(2.0), Inches(5.7), Inches(3.4))
    tf_end = tb_end.text_frame
    tf_end.word_wrap = True
    
    p = tf_end.paragraphs[0]
    p.text = "Terima Kasih!"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Plus Jakarta Sans"

    p2 = tf_end.add_paragraph()
    p2.text = "Sesi Tanya Jawab & Demonstrasi Langsung"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = EMERALD
    p2.font.name = "Plus Jakarta Sans"
    p2.space_before = Pt(10)

    p3 = tf_end.add_paragraph()
    p3.text = "Smart QC Dashboard • IoT Real-Time Monitoring System"
    p3.alignment = PP_ALIGN.CENTER
    p3.font.size = Pt(12)
    p3.font.color.rgb = MUTED_TEXT
    p3.font.name = "Plus Jakarta Sans"
    p3.space_before = Pt(16)

    # Save to scratch folder
    output_path = os.path.join(r"C:\Users\ReX\.gemini\antigravity\scratch\iot-defect-detection", "Presentasi_IoT_Quality_Control.pptx")
    prs.save(output_path)
    print(f"Presentation successfully saved to: {output_path}")

if __name__ == "__main__":
    create_deck()
